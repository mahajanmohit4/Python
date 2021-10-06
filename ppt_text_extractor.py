import pptx

def extract_text_from_pptx(pptx_path):
  prs = pptx.Presentation(pptx_path)
  description = ""
  for slide in prs.slides:
    for shape in slide.shapes:
      if hasattr(shape, "text"):
        description += shape.text
      else:
        description += shape.alt_text
  return description

if __name__ == "__main__":
  pptx_path = "test1.pptx"
  desc = extract_text_from_pptx(pptx_path)
  print (desc)
  
